from datetime import datetime
from io import BytesIO
import logging
import os
import pickle
from pathlib import Path
import re

import nltk
import numpy as np
from docx import Document
from fastapi import Depends, FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from pypdf import PdfReader
from pydantic import BaseModel
from pymongo import ReturnDocument
from nltk.tokenize import sent_tokenize
from pptx import Presentation

from backend.mongo import (
    ensure_collections_and_indexes,
    ensure_default_admin,
    scan_logs_collection,
    users_collection,
)
from backend.security import get_current_user
from features.feature_extractor import build_features_with_chunk_context, warmup_inference_stack
from router.admin import admin_router
from router.auth import auth_router

app = FastAPI(title="Turnitin-Style AI Detector")
logger = logging.getLogger("uvicorn.error")
CHUNK_SENTENCE_SIZE = int(os.getenv("CHUNK_SENTENCE_SIZE", "15"))
MODEL_WARMUP = os.getenv("MODEL_WARMUP", "0").strip().lower() in {"1", "true", "yes", "on"}
MAX_TEXT_CHARS = int(os.getenv("MAX_TEXT_CHARS", "300000"))
MAX_UPLOAD_BYTES = int(os.getenv("MAX_UPLOAD_BYTES", str(20 * 1024 * 1024)))

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "https://ai-checker-blue.vercel.app",
        "http://127.0.0.1:5500",
        "http://localhost:5500",
        "http://127.0.0.1:5501",
        "http://localhost:5501",
        "http://127.0.0.1:3000",
        "http://localhost:3000",
        "http://127.0.0.1:5173",
        "http://localhost:5173",
    ],
    allow_origin_regex=r"https://.*\.vercel\.app|http://(localhost|127\.0\.0\.1|192\.168\.\d{1,3}\.\d{1,3}|10\.\d{1,3}\.\d{1,3}\.\d{1,3}|172\.(1[6-9]|2\d|3[0-1])\.\d{1,3}\.\d{1,3})(:\d+)?$",
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.options("/{path:path}")
async def options_handler(path: str):
    return {"status": "ok"}


@app.get("/")
def health():
    return {"status": "running"}


app.include_router(auth_router)
app.include_router(admin_router)

ensure_collections_and_indexes()
ensure_default_admin()

try:
    nltk.data.find("tokenizers/punkt")
except LookupError:
    nltk.download("punkt", quiet=True)

PROJECT_ROOT = Path(__file__).resolve().parent
MODEL_PATH = PROJECT_ROOT / "models" / "xgb_model_.pkl"
if not MODEL_PATH.exists():
    raise FileNotFoundError(f"Model not found at {MODEL_PATH}")

with MODEL_PATH.open("rb") as f:
    model = pickle.load(f)


@app.on_event("startup")
def warmup_models():
    # On small instances (e.g. 512Mi), eager warmup can OOM. Keep it opt-in.
    if MODEL_WARMUP:
        warmup_inference_stack()


class TextInput(BaseModel):
    text: str


def normalize_extracted_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n").replace("\xa0", " ")
    # Merge PDF-style hyphenated wraps: "exam-\nple" -> "example"
    text = re.sub(r"([A-Za-z])-\n([A-Za-z])", r"\1\2", text)

    lines = text.split("\n")
    cleaned = []
    blank_run = 0

    for raw in lines:
        line = re.sub(r"[ \t]+", " ", raw).strip()
        if not line:
            blank_run += 1
            if blank_run <= 1:
                cleaned.append("")
            continue
        blank_run = 0
        cleaned.append(line)

    return "\n".join(cleaned).strip()


def extract_docx_text(content: bytes) -> str:
    document = Document(BytesIO(content))
    parts = []

    for p in document.paragraphs:
        txt = p.text.strip()
        if not txt:
            continue

        style_name = (p.style.name or "").lower() if p.style else ""
        if style_name.startswith("heading"):
            parts.append(txt.upper())
            parts.append("")
            continue

        # Keep simple bullet/numbered intent visible in plain text output.
        if style_name.startswith("list"):
            parts.append(f"- {txt}")
        else:
            parts.append(txt)

    # Include table content (previously lost).
    for table in document.tables:
        parts.append("")
        for row in table.rows:
            cells = [re.sub(r"\s+", " ", c.text).strip() for c in row.cells]
            cells = [c for c in cells if c]
            if cells:
                parts.append(" | ".join(cells))
        parts.append("")

    return "\n".join(parts)


def extract_pptx_text(content: bytes) -> str:
    try:
        presentation = Presentation(BytesIO(content))
    except Exception:
        raise HTTPException(
            status_code=400,
            detail="Unable to read .ppt file. Please convert it to .pptx and retry.",
        )

    lines = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        lines.append(f"Slide {slide_idx}")
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or shape.text_frame is None:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip() if para.runs else (para.text or "").strip()
                if not text:
                    continue
                level = getattr(para, "level", 0) or 0
                indent = "  " * min(level, 4)
                bullet = "- " if level >= 0 else ""
                lines.append(f"{indent}{bullet}{text}")
        lines.append("")

    return "\n".join(lines)


def extract_pdf_text(content: bytes) -> str:
    reader = PdfReader(BytesIO(content))
    pages = []
    for idx, page in enumerate(reader.pages, start=1):
        try:
            page_text = page.extract_text(extraction_mode="layout") or ""
        except TypeError:
            page_text = page.extract_text() or ""
        if page_text.strip():
            pages.append(page_text.strip())
        else:
            pages.append(f"[Page {idx}: no readable text]")
    return "\n\n".join(pages)


def classify_turnitin(ai_percent, human_percent, polish_percent):
    if ai_percent >= 20 and ai_percent > human_percent:
        return "AI"

    if polish_percent >= 95:
        return "AI"

    if polish_percent < 90:
        return "AI" if ai_percent > human_percent else "Human"

    if human_percent >= ai_percent + 10:
        return "Human"

    if ai_percent >= 12:
        return "AI"

    return "Human"


def consume_user_token(current_user):
    user_doc = users_collection.find_one_and_update(
        {"_id": current_user["_id"], "tokens": {"$gt": 0}},
        {"$inc": {"tokens": -1}},
        return_document=ReturnDocument.BEFORE,
    )

    if not user_doc:
        existing = users_collection.find_one({"_id": current_user["_id"]})
        if not existing:
            raise HTTPException(status_code=403, detail="User Not Found")
        raise HTTPException(status_code=402, detail="TOKEN_FINISHED")

    return user_doc.get("tokens", 0)


def run_prediction(text: str, user_id: str, tokens_before: int):
    text = text.strip()
    if not text:
        raise HTTPException(status_code=400, detail="Text is empty")
    if len(text) > MAX_TEXT_CHARS:
        raise HTTPException(
            status_code=413,
            detail=f"Text too large. Maximum allowed characters: {MAX_TEXT_CHARS}",
        )

    sentences = sent_tokenize(text)
    if not sentences:
        raise HTTPException(status_code=400, detail="No sentences found")

    original_sentence_count = len(sentences)
    results = []
    total_ai = 0
    total_human = 0

    feature_rows = []
    sentence_order = []

    for i in range(0, len(sentences), CHUNK_SENTENCE_SIZE):
        chunk_sentences = sentences[i:i + CHUNK_SENTENCE_SIZE]
        chunk_text = " ".join(chunk_sentences)

        for sent in chunk_sentences:
            feature_rows.append(build_features_with_chunk_context(sent, chunk_text))
            sentence_order.append(sent)

    probs_batch = model.predict_proba(np.vstack(feature_rows))

    for sent, probs in zip(sentence_order, probs_batch):

        human_p = float(probs[0] * 100)
        ai_p = float(probs[1] * 100)
        polish_p = float(probs[2] * 100)

        total_ai += ai_p
        total_human += human_p

        final_label = classify_turnitin(ai_p, human_p, polish_p)
        results.append(
            {
                "sentence": sent,
                "human_probability": round(human_p, 2),
                "ai_probability": round(ai_p, 2),
                "over_polished_probability": round(polish_p, 2),
                "final_label": final_label,
            }
        )

    avg_ai = total_ai / len(sentences)
    avg_human = total_human / len(sentences)
    final_doc_label = "AI" if avg_ai > avg_human else "Human"

    scan_logs_collection.insert_one(
        {
            "uid": user_id,
            "scanned_text": text,
            "result": final_doc_label,
            "ai_percent": round(avg_ai, 2),
            "human_percent": round(avg_human, 2),
            "timestamp": datetime.utcnow(),
        }
    )

    return {
        "tokens_left": max(tokens_before - 1, 0),
        "overall_human_probability": round(avg_human, 2),
        "overall_ai_probability": round(avg_ai, 2),
        "final_document_label": final_doc_label,
        "sentences": results,
        "sentences_processed": len(sentences),
        "sentences_received": original_sentence_count,
    }


def extract_text_from_upload(filename: str, content: bytes) -> str:
    ext = os.path.splitext(filename or "")[1].lower()

    if ext == ".txt":
        try:
            raw = content.decode("utf-8")
        except UnicodeDecodeError:
            raw = content.decode("latin-1", errors="ignore")
        return normalize_extracted_text(raw)

    if ext == ".pdf":
        return normalize_extracted_text(extract_pdf_text(content))

    if ext in {".docx", ".word"}:
        return normalize_extracted_text(extract_docx_text(content))

    if ext in {".ppt", ".pptx"}:
        return normalize_extracted_text(extract_pptx_text(content))

    if ext == ".doc":
        raise HTTPException(
            status_code=400,
            detail="Legacy .doc is not supported directly. Please convert to .docx.",
        )

    raise HTTPException(
        status_code=400,
        detail="Unsupported file type. Use txt, pdf, docx, pptx, or ppt.",
    )


@app.post("/predict")
def predict(data: TextInput, current_user=Depends(get_current_user)):
    user_id = str(current_user["_id"])
    text = data.text.strip()
    tokens_before = consume_user_token(current_user)
    return run_prediction(text=text, user_id=user_id, tokens_before=tokens_before)


@app.post("/predict-file")
async def predict_file(file: UploadFile = File(...), current_user=Depends(get_current_user)):
    user_id = str(current_user["_id"])

    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")
    if len(content) > MAX_UPLOAD_BYTES:
        raise HTTPException(
            status_code=413,
            detail=f"File too large. Maximum allowed size: {MAX_UPLOAD_BYTES} bytes",
        )

    extracted_text = extract_text_from_upload(file.filename or "", content).strip()
    if not extracted_text:
        raise HTTPException(status_code=400, detail="No readable text found in file")

    tokens_before = consume_user_token(current_user)
    return run_prediction(text=extracted_text, user_id=user_id, tokens_before=tokens_before)


@app.post("/extract-file")
async def extract_file(file: UploadFile = File(...), _current_user=Depends(get_current_user)):
    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")
    if len(content) > MAX_UPLOAD_BYTES:
        raise HTTPException(
            status_code=413,
            detail=f"File too large. Maximum allowed size: {MAX_UPLOAD_BYTES} bytes",
        )

    extracted_text = extract_text_from_upload(file.filename or "", content).strip()
    if not extracted_text:
        raise HTTPException(status_code=400, detail="No readable text found in file")

    return {
        "filename": file.filename or "",
        "text": extracted_text,
        "characters": len(extracted_text),
    }


@app.get("/my-history")
def get_my_history(current_user=Depends(get_current_user)):
    user_id = str(current_user["_id"])
    logs = scan_logs_collection.find({"uid": user_id}).sort("timestamp", -1).limit(100)

    return [
        {
            "id": str(log["_id"]),
            "scanned_text": log.get("scanned_text", ""),
            "result": log.get("result"),
            "ai_percent": log.get("ai_percent", 0),
            "human_percent": log.get("human_percent", 0),
            "timestamp": log.get("timestamp").isoformat() if log.get("timestamp") else None,
        }
        for log in logs
    ]
